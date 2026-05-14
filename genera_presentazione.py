"""
Genera la presentazione PowerPoint per il progetto Trasporti AUSL Umbria 1.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as ns
from lxml import etree

# ── PALETTE ─────────────────────────────────────────────────
C_DARK_BG    = RGBColor(0x0F, 0x20, 0x44)   # navy scuro
C_MID_BG     = RGBColor(0x1A, 0x3A, 0x6E)   # blu medio
C_ACCENT_B   = RGBColor(0x25, 0x63, 0xC7)   # blu accento
C_GREEN      = RGBColor(0x34, 0xD3, 0x99)   # verde acqua
C_BLUE       = RGBColor(0x60, 0xA5, 0xFA)   # azzurro chiaro
C_AMBER      = RGBColor(0xF5, 0x9E, 0x0B)   # ambra
C_RED        = RGBColor(0xF8, 0x71, 0x71)   # rosso chiaro
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY       = RGBColor(0x94, 0xA3, 0xB8)
C_DARK_GRAY  = RGBColor(0x47, 0x55, 0x69)
C_PURPLE     = RGBColor(0xA7, 0x8B, 0xFA)
C_TEAL       = RGBColor(0x2D, 0xD4, 0xBF)

# Slide size: widescreen 16:9
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completamente vuoto


# ── HELPERS ─────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txb


def add_textbox_mp(slide, x, y, w, h, paras, wrap=True):
    """paras = list of (text, size, bold, color, align, italic)"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    first = True
    for (text, size, bold, color, align, italic) in paras:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
    return txb


def bg_gradient(slide, color1, color2=None):
    """Full-slide background rectangle."""
    r = add_rect(slide, 0, 0, W, H, color1)
    return r


def accent_stripe(slide, color=None):
    """Left vertical stripe."""
    c = color or C_GREEN
    add_rect(slide, 0, 0, Inches(0.15), H, c)


def top_bar(slide):
    """Top horizontal color bar."""
    add_rect(slide, 0, 0, W, Inches(0.09), C_BLUE)


def slide_label(slide, text):
    add_text(slide, text, Inches(0.35), Inches(0.18), Inches(10), Inches(0.35),
             font_size=9, bold=True, color=C_BLUE)


def slide_number(slide, num):
    add_text(slide, f"{num}/10", W - Inches(1.2), H - Inches(0.4), Inches(1.0), Inches(0.35),
             font_size=9, color=C_DARK_GRAY, align=PP_ALIGN.RIGHT)


def rule_line(slide, y):
    ln = add_rect(slide, Inches(0.35), y, W - Inches(0.7), Inches(0.02), C_BLUE)
    return ln


# ═══════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_gradient(s, C_DARK_BG)
accent_stripe(s, C_GREEN)
top_bar(s)

add_text(s, "RICERCA OPERATIVA & INNOVAZIONE GESTIONALE",
         Inches(0.45), Inches(0.25), Inches(12), Inches(0.4),
         font_size=9, bold=True, color=C_BLUE)

add_text(s, "Ottimizzazione dei Trasporti Sanitari",
         Inches(0.45), Inches(1.05), Inches(12), Inches(0.8),
         font_size=36, bold=True, color=C_WHITE)

add_text(s, "AUSL Umbria 1",
         Inches(0.45), Inches(1.85), Inches(12), Inches(0.7),
         font_size=44, bold=True, color=C_GREEN)

add_textbox_mp(s, Inches(0.45), Inches(2.9), Inches(9), Inches(1.4), [
    ("Un sistema di Linear Programming a 3 livelli per la gestione centralizzata", 15, False, C_GRAY, PP_ALIGN.LEFT, False),
    ("dei trasporti secondari non emergenziali — dalla frammentazione alla Centrale Unica Trasporti (CUT).", 15, False, C_GRAY, PP_ALIGN.LEFT, False),
])

rule_line(s, Inches(5.3))

add_text(s, "AUSL Umbria 1   |   Maggio 2026",
         Inches(0.45), Inches(5.45), Inches(6), Inches(0.35),
         font_size=11, color=C_DARK_GRAY)

# Badge
badge = add_rect(s, Inches(10.3), Inches(5.35), Inches(2.6), Inches(0.55),
                 RGBColor(0x1A, 0x3A, 0x6E))
add_text(s, "RISERVATO – USO INTERNO",
         Inches(10.35), Inches(5.42), Inches(2.5), Inches(0.4),
         font_size=9, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

slide_number(s, 1)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 – IL PROBLEMA OGGI
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_gradient(s, RGBColor(0x4C, 0x05, 0x19))
accent_stripe(s, C_RED)
top_bar(s)
slide_label(s, "CONTESTO – SITUAZIONE ATTUALE")

add_text(s, "La gestione attuale: frammentata e inefficiente",
         Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.6),
         font_size=26, bold=True, color=C_WHITE)

def kpi_card(slide, x, y, w, h, label, value, sub, val_color=None):
    add_rect(slide, x, y, w, h, RGBColor(0x1A, 0x1A, 0x2E))
    add_rect(slide, x, y, Inches(0.06), h, val_color or C_RED)
    add_text(slide, label.upper(), x + Inches(0.12), y + Inches(0.1), w - Inches(0.15), Inches(0.3),
             font_size=8, bold=True, color=C_GRAY)
    add_text(slide, value, x + Inches(0.12), y + Inches(0.35), w - Inches(0.15), Inches(0.5),
             font_size=30, bold=True, color=val_color or C_RED)
    add_text(slide, sub, x + Inches(0.12), y + Inches(0.8), w - Inches(0.15), Inches(0.6),
             font_size=9, color=C_GRAY, wrap=True)

cw = Inches(3.0); ch = Inches(1.6); gap = Inches(0.18)
row1_y = Inches(1.3); row2_y = Inches(3.05)

kpi_card(s, Inches(0.35), row1_y, cw, ch, "Vettori attivi", "24+", "Associazioni\nCRI, Misericordie, ANPAS", C_RED)
kpi_card(s, Inches(0.35) + cw + gap, row1_y, cw, ch, "Convenzioni diverse", "24", "Tariffe eterogenee\nnessuno standard", C_RED)
kpi_card(s, Inches(0.35), row2_y, cw, ch, "Fuori territorio", "15%", "dei trasporti usa\nvettori non locali", C_AMBER)
kpi_card(s, Inches(0.35) + cw + gap, row2_y, cw, ch, "Richieste non evase", "2%", "pazienti senza\nvettore disponibile", C_RED)

# Right panel
rx = Inches(6.6); rw = Inches(6.5)
add_rect(s, rx, Inches(1.25), rw, Inches(5.5), RGBColor(0x1A, 0x1A, 0x2E))
add_rect(s, rx, Inches(1.25), Inches(0.06), Inches(5.5), C_RED)

add_text(s, "⚠  PRINCIPALI CRITICITÀ", rx + Inches(0.15), Inches(1.35), rw - Inches(0.2), Inches(0.35),
         font_size=9, bold=True, color=C_GRAY)

criticita = [
    ("Coordinamento manuale", "Ogni richiesta trattata caso per caso, senza ottimizzazione sistemica."),
    ("Assenza di dati granulari", "Si conoscono i totali annuali, non le singole tratte percorse."),
    ("Costi nascosti", "Vettori fuori territorio, duplicazioni e viaggi a vuoto non sono tracciati."),
    ("Dialisi non ottimizzata", "45 pazienti con rotte ricorrenti gestiti senza consolidamento veicoli."),
]
cy = Inches(1.85)
for title, desc in criticita:
    add_text(s, f"▸  {title}", rx + Inches(0.15), cy, rw - Inches(0.2), Inches(0.28),
             font_size=11, bold=True, color=C_WHITE)
    add_text(s, f"    {desc}", rx + Inches(0.15), cy + Inches(0.27), rw - Inches(0.2), Inches(0.35),
             font_size=10, color=C_GRAY)
    cy += Inches(0.82)

slide_number(s, 2)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 – I NUMERI DEL PROBLEMA
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_gradient(s, RGBColor(0x1E, 0x29, 0x3B))
accent_stripe(s, C_BLUE)
top_bar(s)
slide_label(s, "DATI 2025 – VOLUME DEL SERVIZIO")

add_text(s, "Il perimetro: 1,6 milioni di km all'anno",
         Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.55),
         font_size=26, bold=True, color=C_WHITE)

# Top KPI row
def kpi_top(slide, x, y, w, h, num, unit, desc, col):
    add_rect(slide, x, y, w, h, RGBColor(0x1A, 0x1A, 0x2E))
    add_rect(slide, x, y, w, Inches(0.06), col)
    add_text(slide, num,  x + Inches(0.1), y + Inches(0.15), w - Inches(0.2), Inches(0.55),
             font_size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, unit, x + Inches(0.1), y + Inches(0.68), w - Inches(0.2), Inches(0.25),
             font_size=9, color=C_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x + Inches(0.1), y + Inches(0.92), w - Inches(0.2), Inches(0.35),
             font_size=9, color=C_DARK_GRAY, align=PP_ALIGN.CENTER, wrap=True)

kw = Inches(3.1); kh = Inches(1.4); ky = Inches(1.25); kg = Inches(0.15)
kpi_top(s, Inches(0.35),          ky, kw, kh, "1,6M",   "km / anno",        "Volume totale trasporti",                   C_BLUE)
kpi_top(s, Inches(0.35)+kw+kg,    ky, kw, kh, "~17.000","viaggi / anno",     "Ordinari + dialisi",                        C_GREEN)
kpi_top(s, Inches(0.35)+2*(kw+kg),ky, kw, kh, "4.298",  "km² serviti",       "38 comuni, 6 distretti sanitari",           C_AMBER)
kpi_top(s, Inches(0.35)+3*(kw+kg),ky, kw, kh, "~100",   "richieste / giorno","70-80 secondarie + 20 ospedaliere",         C_PURPLE)

# Table
headers = ["Area", "Viaggi 2025", "km 2025", "km / viaggio"]
rows = [
    ("PO Città di Castello + Umbertide", "3.444", "196.123", "57"),
    ("PO Gubbio-Gualdo Tadino",          "2.348", "127.991", "55"),
    ("POU (Assisi / Media Valle / Castiglione)", "2.260", "163.877", "73"),
    ("6 Distretti territoriali (totale)", "8.924", "688.235", "77"),
    ("Dializzati barellati",              "7.007 sedute", "362.987", "52"),
    ("TOTALE COMPLESSIVO",                "~17.000", "~1.635.000", "—"),
]

ty = Inches(2.85); tx = Inches(0.35); tw_total = Inches(12.5)
col_ws = [Inches(5.5), Inches(2.1), Inches(2.5), Inches(2.0)]

# Header row
hx = tx
for i, h_text in enumerate(headers):
    add_rect(s, hx, ty, col_ws[i], Inches(0.38), RGBColor(0x1A, 0x3A, 0x6E))
    add_text(s, h_text.upper(), hx + Inches(0.08), ty + Inches(0.07), col_ws[i] - Inches(0.1), Inches(0.28),
             font_size=8, bold=True, color=C_BLUE)
    hx += col_ws[i]

ty += Inches(0.38)
for ri, row in enumerate(rows):
    rx2 = tx
    bg = RGBColor(0x1A, 0x1A, 0x2E) if ri % 2 == 0 else RGBColor(0x12, 0x12, 0x25)
    is_total = ri == len(rows) - 1
    if is_total:
        bg = RGBColor(0x0F, 0x2D, 0x5C)
    for ci, cell in enumerate(row):
        add_rect(s, rx2, ty, col_ws[ci], Inches(0.42), bg)
        col = C_WHITE if is_total else (C_AMBER if ci == 2 and ri == 4 else C_GRAY)
        add_text(s, cell, rx2 + Inches(0.08), ty + Inches(0.08), col_ws[ci] - Inches(0.1), Inches(0.3),
                 font_size=10, bold=is_total, color=col)
        rx2 += col_ws[ci]
    ty += Inches(0.42)

slide_number(s, 3)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 – SISTEMA LP A 3 LIVELLI
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_gradient(s, RGBColor(0x1A, 0x3A, 0x6E))
accent_stripe(s, C_GREEN)
top_bar(s)
slide_label(s, "SOLUZIONE – ARCHITETTURA DEL SISTEMA")

add_text(s, "Un sistema LP gerarchico a 3 livelli",
         Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.55),
         font_size=26, bold=True, color=C_WHITE)

levels = [
    ("1", "Livello Strategico — Allocazione Territoriale",
     "Transportation Problem (MIP)  ·  Aggiornamento trimestrale",
     "Assegna i 24 vettori alle 9 zone di servizio minimizzando i km attesi.\nGarantisce ridondanza (min 2 vettori/zona) e bilanciamento tra federazioni.",
     "432 var. intere + 216 binarie", "< 1 secondo", RGBColor(0x1D, 0x4E, 0xD8), C_BLUE,
     "15% → <3%\nfuori territorio"),
    ("2", "Livello Tattico — Assegnamento Giornaliero",
     "Assignment Problem con Time Windows (MIP)  ·  Risolto la notte prima",
     "Assegna i 50-70 trasporti pianificati a vettori specifici rispettando tipo\nveicolo, finestre temporali e limiti di capacità giornaliera.",
     "~1.560 binarie + 2.000 seq.", "10–60 secondi", RGBColor(0x06, 0x5F, 0x46), C_GREEN,
     "-10/15% km\npianificati"),
    ("3", "Livello Operativo — Dispatch Real-Time",
     "Euristica greedy + VRP locale  ·  Risposta immediata",
     "Gestisce richieste urgenti e imprevisti. Filtra vettori disponibili, calcola\nscore di ottimalità e assegna in secondi. Integrato con la CUT.",
     "Algoritmo greedy", "< 5 secondi", RGBColor(0x92, 0x40, 0x0E), C_AMBER,
     "2% → <0.5%\nnon evasi"),
]

ly = Inches(1.3)
for (num, title, subtitle, desc, dim, time_str, bg_col, txt_col, result) in levels:
    lh = Inches(1.68)
    # number badge
    add_rect(s, Inches(0.35), ly, Inches(0.6), lh, bg_col)
    add_text(s, num, Inches(0.35), ly + Inches(0.55), Inches(0.6), Inches(0.6),
             font_size=28, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    # body
    body_w = Inches(7.5)
    add_rect(s, Inches(0.95), ly, body_w, lh, RGBColor(0x0A, 0x17, 0x38))
    add_text(s, title, Inches(1.05), ly + Inches(0.1), body_w - Inches(0.15), Inches(0.35),
             font_size=13, bold=True, color=C_WHITE)
    add_text(s, subtitle, Inches(1.05), ly + Inches(0.42), body_w - Inches(0.15), Inches(0.28),
             font_size=9, color=txt_col)
    add_text(s, desc, Inches(1.05), ly + Inches(0.72), body_w - Inches(0.15), Inches(0.82),
             font_size=10, color=C_GRAY, wrap=True)
    # right panel – metrics
    rp_x = Inches(8.45); rp_w = Inches(2.3)
    add_rect(s, rp_x, ly, rp_w, lh, RGBColor(0x07, 0x10, 0x28))
    add_text(s, "DIMENSIONE", rp_x + Inches(0.1), ly + Inches(0.08), rp_w - Inches(0.15), Inches(0.22),
             font_size=7, bold=True, color=C_DARK_GRAY)
    add_text(s, dim, rp_x + Inches(0.1), ly + Inches(0.28), rp_w - Inches(0.15), Inches(0.35),
             font_size=9, bold=True, color=C_WHITE)
    add_text(s, "TEMPO SOLUZIONE", rp_x + Inches(0.1), ly + Inches(0.65), rp_w - Inches(0.15), Inches(0.22),
             font_size=7, bold=True, color=C_DARK_GRAY)
    add_text(s, time_str, rp_x + Inches(0.1), ly + Inches(0.85), rp_w - Inches(0.15), Inches(0.35),
             font_size=12, bold=True, color=C_GREEN)
    # result panel
    rs_x = Inches(10.75); rs_w = Inches(2.2)
    add_rect(s, rs_x, ly, rs_w, lh, bg_col)
    add_text(s, "IMPATTO ATTESO", rs_x + Inches(0.1), ly + Inches(0.12), rs_w - Inches(0.15), Inches(0.25),
             font_size=8, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    add_text(s, result, rs_x + Inches(0.1), ly + Inches(0.45), rs_w - Inches(0.15), Inches(0.85),
             font_size=14, bold=True, color=C_WHITE, wrap=True)
    ly += lh + Inches(0.14)

slide_number(s, 4)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 – QUICK WIN DIALISI
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_gradient(s, RGBColor(0x0F, 0x4C, 0x5C))
accent_stripe(s, C_TEAL)
top_bar(s)
slide_label(s, "PRIORITÀ – QUICK WIN IMMEDIATO")

add_text(s, "Iniziare dalla dialisi: il problema più strutturato",
         Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.55),
         font_size=26, bold=True, color=C_WHITE)

# Left: reasons
lx = Inches(0.35); lw = Inches(6.2)
add_rect(s, lx, Inches(1.25), lw, Inches(5.7), RGBColor(0x0A, 0x2A, 0x33))
add_rect(s, lx, Inches(1.25), Inches(0.06), Inches(5.7), C_TEAL)

add_text(s, "PERCHÉ LA DIALISI?", lx + Inches(0.15), Inches(1.35), lw - Inches(0.2), Inches(0.3),
         font_size=9, bold=True, color=C_GRAY)

reasons = [
    ("Stessi 45 pazienti ogni settimana", "Rotte ricorrenti 3 volte/settimana — il problema è predittibile al 100%."),
    ("7.007 sedute/anno", "Alto volume, alto impatto sull'ottimizzazione. È il sotto-problema più rilevante."),
    ("Consolidamento immediato", "2 pazienti vicini sullo stesso veicolo dimezza i km di quella rotta."),
    ("Periodic VRP", "Il modello matematico Periodic Vehicle Routing Problem è ideale per rotte ricorrenti."),
]
ry = Inches(1.8)
for title, desc in reasons:
    add_text(s, f"▸  {title}", lx + Inches(0.15), ry, lw - Inches(0.2), Inches(0.3),
             font_size=12, bold=True, color=C_WHITE)
    add_text(s, f"    {desc}", lx + Inches(0.15), ry + Inches(0.3), lw - Inches(0.2), Inches(0.4),
             font_size=10, color=C_GRAY, wrap=True)
    ry += Inches(0.9)

# Right: KPI cards
rx2 = Inches(6.8); rw2 = Inches(6.15)
cards_d = [
    ("Volume km dialisi (2025)",  "362.987 km/anno", "Perugino 211K · Alto Chiascio 82K · Alto Tevere 70K", C_TEAL),
    ("Risparmio stimato (−15/20%)", "54.000 – 72.000 km", "km/anno risparmiati con la sola ottimizzazione dialisi", C_GREEN),
    ("Equivalente economico",      "50.000 – 90.000 €", "risparmio annuo dalla dialisi (≈ 1 €/km risparmiato)", C_AMBER),
]
cy2 = Inches(1.25); ch2 = Inches(1.7)
for label, val, sub, col in cards_d:
    add_rect(s, rx2, cy2, rw2, ch2, RGBColor(0x0A, 0x2A, 0x33))
    add_rect(s, rx2, cy2, Inches(0.06), ch2, col)
    add_text(s, label.upper(), rx2 + Inches(0.15), cy2 + Inches(0.1), rw2 - Inches(0.2), Inches(0.25),
             font_size=8, bold=True, color=C_GRAY)
    add_text(s, val, rx2 + Inches(0.15), cy2 + Inches(0.38), rw2 - Inches(0.2), Inches(0.55),
             font_size=22, bold=True, color=col)
    add_text(s, sub, rx2 + Inches(0.15), cy2 + Inches(0.95), rw2 - Inches(0.2), Inches(0.55),
             font_size=10, color=C_GRAY, wrap=True)
    cy2 += ch2 + Inches(0.18)

# Quote
add_rect(s, Inches(0.35), Inches(6.35), Inches(12.5), Inches(0.82), RGBColor(0x07, 0x1A, 0x20))
add_rect(s, Inches(0.35), Inches(6.35), Inches(0.06), Inches(0.82), C_AMBER)
add_text(s, "\"La dialisi è il quick win più importante: stesse rotte, stessi pazienti ogni settimana."
            " Basta un foglio con i 45 pazienti e le loro sedi per avviare l'ottimizzazione.\"",
         Inches(0.55), Inches(6.45), Inches(12.2), Inches(0.65),
         font_size=11, italic=True, color=RGBColor(0xFD, 0xE6, 0x8A))

slide_number(s, 5)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 – IMPATTO ATTESO / ROI
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_gradient(s, RGBColor(0x1B, 0x43, 0x32))
accent_stripe(s, C_GREEN)
top_bar(s)
slide_label(s, "RISULTATI ATTESI – ROI DEL PROGETTO")

add_text(s, "Impatto quantificato: 230.000 – 350.000 €/anno",
         Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.55),
         font_size=26, bold=True, color=C_WHITE)

# Arrow cards
arrow_cards = [
    ("km annui totali",          "1.539.000 km",    "~1.307.000 km", "−232.000 km (−15%)"),
    ("Fuori territorio",         "15% dei trasporti","< 3%",          "−2.000+ viaggi extra"),
    ("Richieste non evase",      "2% non servite",   "< 0,5%",        "+255 pazienti serviti"),
    ("Tipi di contratto",        "24 eterogenei",    "3 tariffe unif.","1 per federazione"),
]
aw = Inches(3.0); ah = Inches(1.6); ax = Inches(0.35); ay = Inches(1.28); ag = Inches(0.14)
for label, before, after, saving in arrow_cards:
    add_rect(s, ax, ay, aw, ah, RGBColor(0x0A, 0x20, 0x18))
    add_text(s, before, ax + Inches(0.1), ay + Inches(0.1), aw - Inches(0.15), Inches(0.3),
             font_size=10, color=C_RED, bold=True)
    add_text(s, "↓", ax + Inches(0.1), ay + Inches(0.42), aw - Inches(0.15), Inches(0.3),
             font_size=18, color=C_DARK_GRAY)
    add_text(s, after, ax + Inches(0.1), ay + Inches(0.72), aw - Inches(0.15), Inches(0.38),
             font_size=18, bold=True, color=C_GREEN)
    add_text(s, label, ax + Inches(0.1), ay + Inches(1.1), aw - Inches(0.15), Inches(0.25),
             font_size=9, color=C_GRAY)
    add_text(s, saving, ax + Inches(0.1), ay + Inches(1.32), aw - Inches(0.15), Inches(0.22),
             font_size=9, bold=True, color=C_AMBER)
    ax += aw + ag

# Table
headers2 = ["Metrica", "Stato attuale", "Obiettivo", "Risparmio"]
rows2 = [
    ("km totali / anno",        "1.539.000",          "~1.307.000",     "−232.000 km (−15%)"),
    ("km dialisi / anno",       "362.987",             "~300.000",       "−63.000 km (−17%)"),
    ("Costi stimati",           "~2,3–3,5 M€/anno",   "~2,0–3,0 M€",    "−230.000 / −350.000 €"),
    ("Risposta urgenze",        "Manuale, variabile",  "Auto, <30 sec",  "Qualità del servizio"),
]

ty2 = Inches(3.05); tx2 = Inches(0.35)
col_ws2 = [Inches(4.2), Inches(2.8), Inches(2.8), Inches(2.6)]

hx2 = tx2
for h_text in headers2:
    add_rect(s, hx2, ty2, col_ws2[headers2.index(h_text)], Inches(0.38), RGBColor(0x0A, 0x2A, 0x18))
    add_text(s, h_text.upper(), hx2 + Inches(0.08), ty2 + Inches(0.07), col_ws2[headers2.index(h_text)] - Inches(0.1), Inches(0.28),
             font_size=8, bold=True, color=C_GREEN)
    hx2 += col_ws2[headers2.index(h_text)]

ty2 += Inches(0.38)
for ri2, row in enumerate(rows2):
    rx3 = tx2
    bg2 = RGBColor(0x0A, 0x20, 0x18) if ri2 % 2 == 0 else RGBColor(0x07, 0x16, 0x11)
    for ci2, cell in enumerate(row):
        add_rect(s, rx3, ty2, col_ws2[ci2], Inches(0.42), bg2)
        col2 = C_GREEN if ci2 == 3 else (C_WHITE if ci2 == 0 else C_GRAY)
        add_text(s, cell, rx3 + Inches(0.08), ty2 + Inches(0.09), col_ws2[ci2] - Inches(0.1), Inches(0.3),
                 font_size=10, bold=(ci2 == 3), color=col2)
        rx3 += col_ws2[ci2]
    ty2 += Inches(0.42)

# Big number bottom right
add_text(s, "230K – 350K €", Inches(8.5), Inches(5.8), Inches(4.5), Inches(0.8),
         font_size=32, bold=True, color=C_GREEN, align=PP_ALIGN.RIGHT)
add_text(s, "risparmio annuo stimato", Inches(8.5), Inches(6.5), Inches(4.5), Inches(0.35),
         font_size=12, color=C_GRAY, align=PP_ALIGN.RIGHT)

slide_number(s, 6)


# ═══════════════════════════════════════════════════════════
# SLIDE 7 – ROADMAP
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_gradient(s, RGBColor(0x0C, 0x14, 0x45))
accent_stripe(s, C_BLUE)
top_bar(s)
slide_label(s, "PIANO DI LAVORO – 12 MESI")

add_text(s, "Roadmap di implementazione: 4 fasi verso la CUT",
         Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.55),
         font_size=26, bold=True, color=C_WHITE)

phases_data = [
    ("FASE 0", "Raccolta Dati",          "Mesi 1–2",  C_BLUE,
     "▸ Tracciare ogni singola richiesta nel SAT\n▸ Inventario 24 vettori (flotta, orari, basi)\n▸ Matrice distanze stradali reali (ORS)\n\nOutput: database operativo pronto",
     "Nessuna\nottimizzazione\nora", C_BLUE),
    ("FASE 1", "Modello Strategico",     "Mesi 3–4",  C_GREEN,
     "▸ Allocazione territoriale trimestrale\n▸ Assegnazione formale vettori alle zone\n▸ Riduzione viaggi fuori territorio\n\nOutput: −60/70% fuori territorio",
     "−60/70%\nfuori\nterritorio", C_GREEN),
    ("FASE 2", "Modello Tattico",        "Mesi 5–7",  C_AMBER,
     "▸ Assegnamento giornaliero ottimizzato\n▸ Start dalle rotte dialisi (Periodic VRP)\n▸ Estensione a tutti i trasporti pianificati\n\nOutput: −10/15% km pianificati",
     "−10/15%\nkm\npianificati", C_AMBER),
    ("FASE 3+4", "CUT + Contratti",      "Mesi 8–12", C_PURPLE,
     "▸ Dispatch real-time integrato con CUT\n▸ Armonizzazione contratti con prezzi ombra LP\n▸ 24 convenzioni → 3 tariffe standard\n\nOutput: CUT operativa Q1 2027",
     "CUT\noperativa\nQ1 2027", C_PURPLE),
]

pw = Inches(3.0); ph = Inches(4.2); px = Inches(0.35); py = Inches(1.25); pg = Inches(0.12)
for (num, title, period, col, desc, result, rcol) in phases_data:
    # top color bar
    add_rect(s, px, py, pw, Inches(0.08), col)
    # body
    add_rect(s, px, py + Inches(0.08), pw, ph - Inches(0.08), RGBColor(0x08, 0x10, 0x2E))
    add_text(s, num, px + Inches(0.12), py + Inches(0.15), pw - Inches(0.2), Inches(0.28),
             font_size=9, bold=True, color=C_DARK_GRAY)
    add_text(s, title, px + Inches(0.12), py + Inches(0.42), pw - Inches(0.2), Inches(0.45),
             font_size=14, bold=True, color=C_WHITE)
    add_text(s, period, px + Inches(0.12), py + Inches(0.86), pw - Inches(0.2), Inches(0.28),
             font_size=10, bold=True, color=col)
    add_text(s, desc, px + Inches(0.12), py + Inches(1.2), pw - Inches(0.2), Inches(2.0),
             font_size=9.5, color=C_GRAY, wrap=True)
    # result box
    add_rect(s, px, py + ph - Inches(1.0), pw, Inches(1.0), col)
    add_text(s, result, px + Inches(0.1), py + ph - Inches(0.95), pw - Inches(0.15), Inches(0.9),
             font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    px += pw + pg

slide_number(s, 7)


# ═══════════════════════════════════════════════════════════
# SLIDE 8 – STACK TECNOLOGICO
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_gradient(s, RGBColor(0x3B, 0x07, 0x64))
accent_stripe(s, C_PURPLE)
top_bar(s)
slide_label(s, "TECNOLOGIA – SOLUZIONI OPEN SOURCE E GRATUITE")

add_text(s, "Stack tecnologico: zero costi di licenza",
         Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.55),
         font_size=26, bold=True, color=C_WHITE)

tech_items = [
    ("🔢  PuLP + CBC Solver",       "Motore LP/MIP — Python, gratuito. Risolve i problemi AUSL Umbria 1 in secondi.",         C_GREEN),
    ("🗺️  OpenRouteService API",    "Matrice distanze stradali reali. Gratuito fino a 2.000 richieste/giorno.",               C_BLUE),
    ("🗄️  PostgreSQL + PostGIS",    "Database relazionale con estensione geografica. Standard adottato in sanità pubblica.",   C_PURPLE),
    ("📊  Streamlit Dashboard",     "Dashboard KPI real-time per operatori CUT. Deployabile in 1 giorno.",                    C_AMBER),
    ("📥  Import CSV da SAT",       "Integrazione non invasiva: nessuna modifica ai sistemi gestionali esistenti.",            C_RED),
    ("🐍  Python 3.x",              "Ecosistema ricco per LP e ML — manutenibile e estendibile internamente.",                 C_TEAL),
]

col1 = tech_items[:3]
col2 = tech_items[3:]
ty3 = Inches(1.3)
for col_items, cx3 in [(col1, Inches(0.35)), (col2, Inches(6.8))]:
    for (name, desc, col3) in col_items:
        add_rect(s, cx3, ty3, Inches(6.15), Inches(1.35), RGBColor(0x18, 0x05, 0x30))
        add_rect(s, cx3, ty3, Inches(0.06), Inches(1.35), col3)
        add_text(s, name, cx3 + Inches(0.15), ty3 + Inches(0.12), Inches(5.8), Inches(0.4),
                 font_size=12, bold=True, color=C_WHITE)
        add_text(s, desc, cx3 + Inches(0.15), ty3 + Inches(0.55), Inches(5.8), Inches(0.65),
                 font_size=10, color=C_GRAY, wrap=True)
        ty3 += Inches(1.5)
    ty3 = Inches(1.3)

# Bottom summary
add_rect(s, Inches(0.35), Inches(5.85), Inches(12.5), Inches(0.85), RGBColor(0x0E, 0x04, 0x1C))
add_rect(s, Inches(0.35), Inches(5.85), Inches(0.06), Inches(0.85), C_GREEN)
add_textbox_mp(s, Inches(0.55), Inches(5.95), Inches(12.2), Inches(0.7), [
    ("Costo infrastruttura: 0 € di licenze + costi server standard ASL  —  ", 12, True, C_GREEN, PP_ALIGN.LEFT, False),
    ("Nessun vendor lock-in  ·  Competenze mantenibili internamente  ·  ROI massimizzato", 12, False, C_GRAY, PP_ALIGN.LEFT, False),
])

slide_number(s, 8)


# ═══════════════════════════════════════════════════════════
# SLIDE 9 – RISCHI E MITIGAZIONI
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_gradient(s, RGBColor(0x45, 0x1A, 0x03))
accent_stripe(s, C_AMBER)
top_bar(s)
slide_label(s, "GESTIONE RISCHI – PRINCIPALI SFIDE")

add_text(s, "Rischi e strategie di mitigazione",
         Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.55),
         font_size=26, bold=True, color=C_WHITE)

risk_rows = [
    ("Qualità dati insufficiente",   "Alta",   "Alto",   "Fase 0 dedicata raccolta dati; validazione su 10 tratte campione vs Google Maps"),
    ("Resistenza vettori volontari", "Media",  "Medio",  "Modello propone, vettore conferma; buffer 10% capacità; coinvolgimento anticipato"),
    ("Integrazione gestionale SAT",  "Media",  "Medio",  "Import manuale CSV come fallback; contatto con fornitore SAT nella Fase 0"),
    ("Territorio montano non ott.",  "Bassa",  "Medio",  "Modello usa distanze stradali reali; realtà geografica ≠ inefficienza"),
    ("Solver tattico troppo lento",  "Bassa",  "Medio",  "CBC risolve in <60 sec; fallback greedy disponibile se necessario"),
]

ty4 = Inches(1.25)
rh4 = Inches(0.95)
headers4 = ["Rischio", "Prob.", "Impatto", "Mitigazione"]
col_ws4 = [Inches(3.5), Inches(1.2), Inches(1.2), Inches(7.2)]
hx4 = Inches(0.35)
for hi4, h4 in enumerate(headers4):
    add_rect(s, hx4, ty4, col_ws4[hi4], Inches(0.42), RGBColor(0x30, 0x10, 0x02))
    add_text(s, h4.upper(), hx4 + Inches(0.08), ty4 + Inches(0.09), col_ws4[hi4] - Inches(0.12), Inches(0.28),
             font_size=8, bold=True, color=C_AMBER)
    hx4 += col_ws4[hi4]

ty4 += Inches(0.42)
prob_colors = {"Alta": C_RED, "Media": C_AMBER, "Bassa": C_GREEN}
for ri4, (risk, prob, impact, mit) in enumerate(risk_rows):
    bg4 = RGBColor(0x20, 0x0B, 0x01) if ri4 % 2 == 0 else RGBColor(0x18, 0x08, 0x01)
    rx4 = Inches(0.35)
    cells4 = [risk, prob, impact, mit]
    for ci4, cell4 in enumerate(cells4):
        add_rect(s, rx4, ty4, col_ws4[ci4], rh4, bg4)
        if ci4 == 1:
            c4 = prob_colors.get(cell4, C_WHITE)
        elif ci4 == 2:
            c4 = C_RED if cell4 == "Alto" else C_AMBER if cell4 == "Medio" else C_GREEN
        elif ci4 == 0:
            c4 = C_WHITE
        else:
            c4 = C_GRAY
        add_text(s, cell4, rx4 + Inches(0.08), ty4 + Inches(0.1), col_ws4[ci4] - Inches(0.12), rh4 - Inches(0.15),
                 font_size=10, bold=(ci4 == 1 or ci4 == 2), color=c4, wrap=True)
        rx4 += col_ws4[ci4]
    ty4 += rh4 + Inches(0.06)

# Principle box
add_rect(s, Inches(0.35), Inches(6.2), Inches(12.5), Inches(0.75), RGBColor(0x1A, 0x09, 0x01))
add_rect(s, Inches(0.35), Inches(6.2), Inches(0.06), Inches(0.75), C_AMBER)
add_text(s, "Principio guida:  Affidabilità > Costo — la funzione obiettivo bilancia costo/km con lo score di affidabilità storica di ogni vettore.",
         Inches(0.55), Inches(6.28), Inches(12.2), Inches(0.55),
         font_size=11, italic=True, color=RGBColor(0xFD, 0xE6, 0x8A))

slide_number(s, 9)


# ═══════════════════════════════════════════════════════════
# SLIDE 10 – CALL TO ACTION
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_gradient(s, C_DARK_BG)
accent_stripe(s, C_GREEN)
top_bar(s)
slide_label(s, "PROSSIMI PASSI – DECISIONI NECESSARIE")

add_text(s, "Cosa serve per partire",
         Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.55),
         font_size=32, bold=True, color=C_WHITE)

# Left: action list
lx2 = Inches(0.35); lw2 = Inches(6.5)
add_rect(s, lx2, Inches(1.25), lw2, Inches(5.0), RGBColor(0x0A, 0x17, 0x38))
add_rect(s, lx2, Inches(1.25), Inches(0.06), Inches(5.0), C_BLUE)
add_text(s, "AZIONI IMMEDIATE (entro 30 giorni)",
         lx2 + Inches(0.15), Inches(1.35), lw2 - Inches(0.2), Inches(0.3),
         font_size=9, bold=True, color=C_BLUE)

actions = [
    "Avviare tracciamento granulare di ogni richiesta di trasporto nel gestionale SAT",
    "Inventariare i 24 vettori: sede, flotta, orari operativi, costo/km contrattuale",
    "Designare un referente tecnico interno (data analyst / IT) per il progetto",
    "Raccogliere l'elenco dei 45 pazienti dializzati con indirizzi e centri dialisi",
    "Approvare il budget di sviluppo (stima: 1 FTE per 12 mesi)",
]
ay2 = Inches(1.8)
for act in actions:
    add_text(s, f"→  {act}", lx2 + Inches(0.15), ay2, lw2 - Inches(0.2), Inches(0.55),
             font_size=11, color=C_GRAY, wrap=True)
    ay2 += Inches(0.72)

# Right: 3 cards
rx4 = Inches(7.05); rw4 = Inches(5.9)
right_cards = [
    ("Primo risultato tangibile", "Settimana 8–10",
     "Ottimizzazione rotte dialisi attiva → risparmio di 54.000–72.000 km/anno verificabile con dati reali.",
     C_GREEN),
    ("Obiettivo strategico", "Q1 2027 – CUT operativa",
     "Centrale Unica Trasporti pienamente operativa presso l'Ospedale della Media Valle del Tevere.",
     C_BLUE),
    ("ROI atteso", "230K – 350K €/anno",
     "Risparmio annuo ricorrente. Payback dell'investimento stimato entro il primo anno di piena operatività.",
     C_AMBER),
]
cy3 = Inches(1.25); ch3 = Inches(1.58)
for (title_c, val_c, desc_c, col_c) in right_cards:
    add_rect(s, rx4, cy3, rw4, ch3, RGBColor(0x0A, 0x17, 0x38))
    add_rect(s, rx4, cy3, Inches(0.06), ch3, col_c)
    add_text(s, title_c.upper(), rx4 + Inches(0.15), cy3 + Inches(0.1), rw4 - Inches(0.2), Inches(0.28),
             font_size=8, bold=True, color=C_GRAY)
    add_text(s, val_c, rx4 + Inches(0.15), cy3 + Inches(0.35), rw4 - Inches(0.2), Inches(0.42),
             font_size=18, bold=True, color=col_c)
    add_text(s, desc_c, rx4 + Inches(0.15), cy3 + Inches(0.8), rw4 - Inches(0.2), Inches(0.65),
             font_size=10, color=C_GRAY, wrap=True)
    cy3 += ch3 + Inches(0.15)

# Footer
rule_line(s, Inches(6.55))
add_text(s, "AUSL Umbria 1   ·   Progetto Centrale Unica Trasporti   ·   Maggio 2026",
         Inches(0.4), Inches(6.7), Inches(7), Inches(0.35),
         font_size=10, color=C_DARK_GRAY)
add_text(s, "danielegrottiuk@gmail.com",
         Inches(9.8), Inches(6.7), Inches(3.1), Inches(0.35),
         font_size=10, color=C_BLUE, align=PP_ALIGN.RIGHT)

slide_number(s, 10)


# ── SAVE ────────────────────────────────────────────────────
output_path = r"c:\Users\danie\Il mio Drive\Lesson\Umbria\Trasporti\ricerca_operativa_umbria\Presentazione_AUSL_Umbria1_CUT.pptx"
prs.save(output_path)
print(f"Salvato: {output_path}")
